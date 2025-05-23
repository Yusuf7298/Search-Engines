import os
import io
from django.core.management.base import BaseCommand
from Search_Engine.models import Document, DocumentChunk
import PyPDF2
import sqlite3

from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd


class Command(BaseCommand):
    help = 'Loads documents from a directory into the database'

    def add_arguments(self, parser):
        parser.add_argument('directory', type=str, help='Directory containing the documents')

    def handle(self, *args, **kwargs):
        directory = kwargs['directory']
        db_path = "db.sqlite3"

        if not os.path.exists(db_path):
            self.stdout.write(self.style.ERROR(f"Database file not found at path: {db_path}"))
            return

        for filename in os.listdir(directory):
            filepath = os.path.join(directory, filename)

            if filename.endswith(".txt"):
                try:
                    with open(filepath, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    self.stdout.write(self.style.WARNING(f"Could not decode {filename} as UTF-8. Trying latin-1..."))
                    try:
                        with open(filepath, 'r', encoding='latin-1') as f:
                            content = f.read()
                    except Exception as e:
                        self.stdout.write(self.style.ERROR(f"Error loading {filename}: {e}"))
                        continue
                except FileNotFoundError as e:
                    self.stdout.write(self.style.ERROR(f"File not found: {filename}. Error: {e}"))
                    continue
                except Exception as e:
                    self.stdout.write(self.style.ERROR(f"Unexpected error reading {filename}: {e}"))
                    continue

            elif filename.endswith(".pdf"):
                try:
                    with open(filepath, 'rb') as f:
                        try:
                            pdf_reader = PyPDF2.PdfReader(f)
                            content = ""
                            for page_num in range(len(pdf_reader.pages)):
                                page = pdf_reader.pages[page_num]
                                content += page.extract_text()
                        except Exception as pdf_err:
                            self.stdout.write(self.style.ERROR(f"Error processing PDF content for {filename}: {pdf_err}"))
                            content = None
                except FileNotFoundError as e:
                    self.stdout.write(self.style.ERROR(f"PDF File not found: {filename}. Error: {e}"))
                    content = None
                except Exception as e:
                    self.stdout.write(self.style.ERROR(f"Error loading PDF {filename}: {e}"))
                    content = None

                if content is None:
                    continue

            else:
                self.stdout.write(self.style.WARNING(f"Skipping {filename}: Unsupported file type"))
                continue

            title = filename.rsplit('.', 1)[0]
            try:
                document = Document.objects.create(title=title, content=content)
                chunk_size = 1000
                for i in range(0, len(content), chunk_size):
                    chunk_content = content[i:i + chunk_size]
                    DocumentChunk.objects.create(document=document, content=chunk_content)

                self.stdout.write(self.style.SUCCESS(f'Successfully loaded document: {filename}'))

            except sqlite3.Error as db_err:
                self.stdout.write(self.style.ERROR(f"Database error creating document/chunks for {filename}: {db_err}"))
            except Exception as e:
                self.stdout.write(self.style.ERROR(f"General error creating document/chunks for {filename}: {e}"))
            content = ""

            try:
                if filename.endswith(".txt"):
                    with open(filepath, 'r', encoding='utf-8') as f:
                        content = f.read()

                elif filename.endswith(".pdf"):
                    with open(filepath, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        content = ''.join([page.extract_text() or "" for page in reader.pages])

                elif filename.endswith(".docx"):
                    doc = DocxDocument(filepath)
                    content = '\n'.join([para.text for para in doc.paragraphs])

                elif filename.endswith(".pptx"):
                    prs = Presentation(filepath)
                    slides = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slides.append(shape.text)
                    content = '\n'.join(slides)

                elif filename.endswith(".csv") or filename.endswith(".xlsx"):
                    df = pd.read_csv(filepath) if filename.endswith(".csv") else pd.read_excel(filepath)
                    content = df.to_string(index=False)

                else:
                    self.stdout.write(self.style.WARNING(f"Skipping {filename}: Unsupported file type"))
                    continue

                # Save to database
                title = filename.rsplit('.', 1)[0]
                document = Document.objects.create(title=title, content=content)

                chunk_size = 1000
                for i in range(0, len(content), chunk_size):
                    chunk = content[i:i + chunk_size]
                    DocumentChunk.objects.create(document=document, content=chunk)

                self.stdout.write(self.style.SUCCESS(f"Successfully loaded: {filename}"))

            except Exception as e:
                self.stdout.write(self.style.ERROR(f"Failed to load {filename}: {e}"))
